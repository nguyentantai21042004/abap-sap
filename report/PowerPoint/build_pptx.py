#!/usr/bin/env python3
"""
Build ZBUG_WS presentation.pptx using python-pptx.
Design: navy (#002060) bg, cyan (#00C7FF) accents, pink (#DE73C9) headings.
Output: report/PowerPoint/output/presentation_py.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color palette ──────────────────────────────────────────────────────────────
BG = RGBColor(0x00, 0x20, 0x60)
BG_DARK = RGBColor(0x00, 0x10, 0x40)
CYAN = RGBColor(0x00, 0xC7, 0xFF)
PINK = RGBColor(0xDE, 0x73, 0xC9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xD8, 0xE8, 0xF4)
MUTED = RGBColor(0xA0, 0xB4, 0xCC)
TH_BG = RGBColor(0x00, 0x88, 0xBB)
TH_TEXT = RGBColor(0x00, 0x20, 0x60)
TD_ODD = RGBColor(0x00, 0x28, 0x72)
TD_EVEN = RGBColor(0x00, 0x1E, 0x58)
CODE_BG = RGBColor(0x00, 0x08, 0x1A)
CODE_FG = RGBColor(0xA8, 0xD8, 0xF0)
NOTE_BG = RGBColor(0x00, 0x24, 0x68)

# ── Slide dimensions ───────────────────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

BAR_W = Inches(0.10)  # left accent bar
LEFT = Inches(0.76)  # content left edge
RIGHT = W - Inches(0.45)
CW = int(RIGHT) - int(LEFT)  # content width in EMU

TITLE_Y = Inches(0.30)
TITLE_H = Inches(0.62)
SEP_Y = Inches(0.93)
SEP_H = Emu(26000)

BODY_Y = Inches(1.07)
BODY_H = H - BODY_Y - Inches(0.48)

FOOT_Y = H - Inches(0.43)
FOOT_H = Inches(0.38)

TOTAL = 16  # total slide count

FONT = "Segoe UI"
MONO = "Courier New"


# ── Low-level helpers ──────────────────────────────────────────────────────────


def _rect(slide, x, y, w, h, color, border=False):
    """Filled rectangle shape, no border by default."""
    x, y, w, h = int(x), int(y), int(w), int(h)
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if border:
        shp.line.color.rgb = color
    else:
        shp.line.fill.background()
    return shp


def _txb(slide, x, y, w, h, wrap=True):
    """Add textbox, return its text_frame."""
    txb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    txb.text_frame.word_wrap = wrap
    return txb.text_frame


def _run(para, text, size=Pt(16), color=LIGHT, bold=False, mono=False, align=None):
    """Append a run to a paragraph with given style."""
    if align is not None:
        para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = MONO if mono else FONT
    return run


def _cell_text(cell, text, size=Pt(13), color=LIGHT, bold=False, align=PP_ALIGN.LEFT):
    """Set text in a table cell with formatting."""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    # Clear any existing runs
    for run in list(p.runs):
        run._r.getparent().remove(run._r)
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = FONT


# ── Slide-level helpers ────────────────────────────────────────────────────────


def _new_slide(prs):
    """Blank navy slide with left accent bar and bottom line."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Left accent bar: cyan top half, pink bottom half
    _rect(slide, 0, 0, BAR_W, H // 2, CYAN)
    _rect(slide, 0, H // 2, BAR_W, H // 2, PINK)

    # Bottom accent line
    _rect(slide, 0, H - Emu(30000), W, Emu(30000), CYAN)
    return slide


def _title(slide, text):
    """Pink heading + cyan separator."""
    tf = _txb(slide, LEFT, TITLE_Y, CW, TITLE_H, wrap=False)
    _run(tf.paragraphs[0], text, size=Pt(27), color=PINK, bold=True)
    _rect(slide, LEFT, SEP_Y, CW, SEP_H, CYAN)


def _footer(slide, text, page_num):
    """Footer note left + page number right."""
    if text:
        tf = _txb(slide, LEFT, FOOT_Y, CW - Inches(1.3), FOOT_H)
        _run(tf.paragraphs[0], text, size=Pt(11.5), color=MUTED)
    tf2 = _txb(slide, W - Inches(1.25), FOOT_Y, Inches(1.15), FOOT_H)
    _run(
        tf2.paragraphs[0],
        f"{page_num} / {TOTAL}",
        size=Pt(11),
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def _note(slide, text, y, height=Inches(0.56)):
    """Cyan-border note box."""
    _rect(slide, LEFT, y, Inches(0.05), height, CYAN)
    _rect(slide, LEFT + Inches(0.05), y, CW - Inches(0.05), height, NOTE_BG)
    tf = _txb(
        slide,
        LEFT + Inches(0.20),
        y + Inches(0.08),
        CW - Inches(0.28),
        height - Inches(0.13),
    )
    _run(tf.paragraphs[0], text, size=Pt(13), color=LIGHT)


def _label(slide, text, y, color=CYAN):
    """Small bold label line."""
    tf = _txb(slide, LEFT, y, CW, Inches(0.45))
    _run(tf.paragraphs[0], text, size=Pt(15.5), color=color, bold=True)


def _code_block(slide, lines, y, height):
    """Dark monospace code block with cyan left bar."""
    _rect(slide, LEFT, y, CW, height, CODE_BG)
    _rect(slide, LEFT, y, Inches(0.05), height, CYAN)
    tf = _txb(
        slide,
        LEFT + Inches(0.18),
        y + Inches(0.12),
        CW - Inches(0.25),
        height - Inches(0.20),
    )
    tf.word_wrap = False
    p = tf.paragraphs[0]
    _run(p, "\n".join(lines), size=Pt(12), color=CODE_FG, mono=True)


def _col_widths(proportions):
    """Convert proportion list to Emu widths summing to CW."""
    w = [int(CW * p) for p in proportions]
    w[-1] += CW - sum(w)
    return w


def _table(slide, headers, rows, proportions, y, height, fs=Pt(13)):
    """Add a fully styled table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    widths = _col_widths(proportions)

    frame = slide.shapes.add_table(n_rows, n_cols, int(LEFT), int(y), CW, int(height))
    tbl = frame.table

    # Set column widths
    for i, w in enumerate(widths):
        tbl.columns[i].width = w

    # Header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TH_BG
        _cell_text(cell, hdr, size=fs, color=TH_TEXT, bold=True)

    # Data rows
    for ri, row in enumerate(rows):
        bg = TD_ODD if ri % 2 == 0 else TD_EVEN
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            # Detect **bold** markers
            s = str(val)
            is_bold = s.startswith("**") and s.endswith("**")
            clean = s[2:-2] if is_bold else s
            _cell_text(
                cell, clean, size=fs, color=WHITE if is_bold else LIGHT, bold=is_bold
            )

    return tbl


# ── Individual slides ──────────────────────────────────────────────────────────


def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_DARK

    # Accent bar
    _rect(slide, 0, 0, Inches(0.12), int(H * 0.55), CYAN)
    _rect(slide, 0, int(H * 0.55), Inches(0.12), int(H * 0.45), PINK)
    _rect(slide, 0, H - Emu(30000), W, Emu(30000), CYAN)

    # Title
    tf1 = _txb(slide, Inches(0.90), Inches(1.80), Inches(11.5), Inches(1.40))
    _run(
        tf1.paragraphs[0],
        "SAP Bug Tracking Management System",
        size=Pt(38),
        color=CYAN,
        bold=True,
    )

    # Separator
    _rect(slide, Inches(0.90), Inches(3.30), Inches(11.0), Emu(42000), CYAN)

    # Subtitle
    tf2 = _txb(slide, Inches(0.90), Inches(3.48), Inches(8.0), Inches(0.80))
    _run(tf2.paragraphs[0], "ZBUG_WS \u2014 v5.0", size=Pt(28), color=PINK, bold=True)

    # Meta line
    tf3 = _txb(slide, Inches(0.90), Inches(4.50), Inches(10.0), Inches(0.60))
    _run(
        tf3.paragraphs[0],
        "Nh\u00f3m ZBUG  \u2502  FPT University Capstone  \u2502  Th\u00e1ng 4 n\u0103m 2026",
        size=Pt(16),
        color=MUTED,
    )

    # Page num
    tf4 = _txb(slide, W - Inches(1.25), FOOT_Y, Inches(1.15), FOOT_H)
    _run(
        tf4.paragraphs[0],
        f"1 / {TOTAL}",
        size=Pt(11),
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def slide_02_team(prs):
    slide = _new_slide(prs)
    _title(slide, "Nh\u00f3m D\u1ef1 \u00c1n")

    headers = [
        "T\u00e0i kho\u1ea3n",
        "H\u1ecd t\u00ean",
        "M\u00e3 SV",
        "Vai tr\u00f2",
        "Tr\u00e1ch nhi\u1ec7m ch\u00ednh",
    ]
    rows = [
        [
            "DEV-237",
            "Nguy\u1ec5n Ng\u1ecdc \u0110\u1ee9c",
            "SE183121",
            "**Leader**",
            "V\u00f2ng \u0111\u1eddi tr\u1ea1ng th\u00e1i, Auto-Assign, Qu\u1ea3n l\u00fd D\u1ef1 \u00e1n",
        ],
        [
            "DEV-089",
            "Nguy\u1ec5n Ho\u00e0ng Anh",
            "SE173545",
            "Developer",
            "Thi\u1ebft k\u1ebf CSDL, logic ABAP l\u00f5i, t\u00e0i li\u1ec7u",
        ],
        [
            "DEV-242",
            "Nguy\u1ec5n H\u00e0 Linh",
            "SS170495",
            "Developer",
            "M\u00e0n h\u00ecnh Bug Detail (0300), FM t\u1ea1o & ghi log",
        ],
        [
            "DEV-061",
            "Nguy\u1ec5n Tr\u1ecdng Hi\u1ebfu",
            "SE180504",
            "Developer",
            "Bug List + Dashboard (0200), Search Engine",
        ],
        [
            "DEV-118",
            "B\u00f9i Anh Kha",
            "SE181730",
            "Tester / QC",
            "Email, Upload b\u1eb1ng ch\u1ee9ng, SmartForms, QA",
        ],
    ]
    _table(
        slide,
        headers,
        rows,
        [0.09, 0.17, 0.10, 0.10, 0.54],
        BODY_Y,
        Inches(3.85),
        fs=Pt(13),
    )
    _footer(
        slide,
        "T\u1ed5ng n\u1ed7 l\u1ef1c: ~35 ng\u00e0y-ng\u01b0\u1eddi  |  H\u1ec7 th\u1ed1ng: SAP S40, Client 324, ABAP 7.70",
        2,
    )


def slide_03_context(prs):
    slide = _new_slide(prs)
    _title(slide, "B\u1ed1i C\u1ea3nh & V\u1ea5n \u0110\u1ec1")

    tf = _txb(slide, LEFT, BODY_Y, CW, Inches(1.60))
    _run(
        tf.paragraphs[0],
        "Tri\u1ec3n khai SAP ERP song song nhi\u1ec1u module (MM, SD, FI, CO\u2026) \u2014 kh\u00f4ng c\u00f3 c\u00f4ng c\u1ee5 theo d\u00f5i l\u1ed7i t\u1eadp trung:",
        size=Pt(15.5),
        color=LIGHT,
    )

    # Problem section
    _label(
        slide, "H\u1eadu qu\u1ea3 th\u1ef1c t\u1ebf:", BODY_Y + Inches(0.55), color=PINK
    )
    problems = [
        "B\u00e1o c\u00e1o l\u1ed7i ph\u00e2n t\u00e1n qua b\u1ea3ng t\u00ednh, email, Zalo",
        "Developer nh\u1eadn ph\u00e2n c\u00f4ng kh\u00f4ng r\u00f5 \u01b0u ti\u00ean hay module",
        "Manager kh\u00f4ng th\u1ea5y ph\u00e2n b\u1ed5 workload theo th\u1eddi gian th\u1ef1c",
        "Kh\u00f4ng c\u00f3 b\u1eb1ng ch\u1ee9ng c\u00f3 c\u1ea5u tr\u00fac v\u1ec1 b\u00e1o c\u00e1o & x\u00e1c nh\u1eadn s\u1eeda l\u1ed7i",
    ]
    tf2 = _txb(slide, LEFT, BODY_Y + Inches(1.02), CW, Inches(1.70))
    first = True
    for prob in problems:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        _run(p, f"  \u2022  {prob}", size=Pt(15), color=LIGHT)

    # Solution box
    sol_y = BODY_Y + Inches(2.90)
    _rect(slide, LEFT, sol_y, CW, Inches(2.35), RGBColor(0x00, 0x22, 0x68))
    _rect(slide, LEFT, sol_y, Inches(0.05), Inches(2.35), CYAN)
    _label(slide, "Gi\u1ea3i ph\u00e1p \u2014 ZBUG_WS:", sol_y + Inches(0.06))
    tf3 = _txb(
        slide,
        LEFT + Inches(0.20),
        sol_y + Inches(0.52),
        CW - Inches(0.25),
        Inches(1.72),
    )
    for i, sol in enumerate(
        [
            "Ch\u1ea1y native tr\u00ean SAP \u2014 kh\u00f4ng c\u1ea7n license b\u1ed5 sung",
            "Truy c\u1eadp qua 1 T-code duy nh\u1ea5t: ZBUG_WS",
            "To\u00e0n b\u1ed9 d\u1eef li\u1ec7u trong b\u1ea3ng ABAP, UI l\u00e0 SAP GUI Dynpro",
        ]
    ):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        _run(p, f"  \u2022  {sol}", size=Pt(15), color=LIGHT)

    _footer(slide, "", 3)


def slide_04_comparison(prs):
    slide = _new_slide(prs)
    _title(slide, "So S\u00e1nh: H\u1ec7 Th\u1ed1ng Tham Chi\u1ebfu vs ZBUG_WS")

    headers = ["T\u00ednh n\u0103ng", "ZPG (Tham chi\u1ebfu)", "ZBUG_WS v5.0"]
    rows = [
        ["Ki\u1ebfn tr\u00fac Module Pool (Type M)", "✗", "✓"],
        [
            "V\u00f2ng \u0111\u1eddi 10 tr\u1ea1ng th\u00e1i + Popup chuy\u1ec3n \u0111\u1ed5i",
            "✗",
            "✓",
        ],
        ["Auto-assign theo SAP module / workload", "✗", "✓"],
        ["Ph\u00e2n quy\u1ec1n t\u1eadp trung qua Function Module", "✗", "✓"],
        ["Audit log \u0111\u1ea7y \u0111\u1ee7 (ZBUG_HISTORY)", "✗", "✓"],
        ["Th\u00f4ng b\u00e1o email qua CL_BCS", "✗", "✓"],
        ["Xu\u1ea5t PDF SmartForm", "✗", "✓"],
        ["Dashboard Header realtime", "✗", "✓"],
        ["Bug Search Engine (0210/0220)", "✗", "✓"],
        ["Module Qu\u1ea3n l\u00fd D\u1ef1 \u00e1n", "✗", "✓"],
    ]
    _table(slide, headers, rows, [0.68, 0.16, 0.16], BODY_Y, Inches(5.55), fs=Pt(13.5))
    _footer(slide, "", 4)


def slide_05_screens(prs):
    slide = _new_slide(prs)
    _title(
        slide,
        "T\u1ed5ng Quan H\u1ec7 Th\u1ed1ng \u2014 B\u1ea3n \u0110\u1ed3 M\u00e0n H\u00ecnh",
    )

    headers = ["M\u00e0n h\u00ecnh", "T\u00ean", "M\u1edbi v5.0"]
    rows = [
        [
            "0410",
            "T\u00ecm ki\u1ebfm D\u1ef1 \u00e1n \u2014 \u0111i\u1ec3m v\u00e0o T-code",
            "✓",
        ],
        ["0400", "Danh s\u00e1ch D\u1ef1 \u00e1n (ALV Grid)", ""],
        ["0200", "Bug List + Dashboard Header realtime", "✓"],
        ["0300", "Bug Detail \u2014 Tab Strip 6 tab", ""],
        ["0370", "Popup Chuy\u1ec3n Tr\u1ea1ng Th\u00e1i (Modal Dialog)", "✓"],
        [
            "0500",
            "Chi ti\u1ebft D\u1ef1 \u00e1n + ph\u00e2n c\u00f4ng ng\u01b0\u1eddi d\u00f9ng",
            "",
        ],
        [
            "0210/0220",
            "Bug Search \u2014 l\u1ecdc \u0111a tr\u01b0\u1eddng + k\u1ebft qu\u1ea3 ALV",
            "✓",
        ],
    ]
    _table(slide, headers, rows, [0.12, 0.76, 0.12], BODY_Y, Inches(4.30), fs=Pt(13.5))
    _footer(
        slide,
        "G\u00f3i: ZBUGTRACK  |  Ch\u01b0\u01a1ng tr\u00ecnh: Z_BUG_WORKSPACE_MP  |  6 ABAP Includes",
        5,
    )


def slide_06_code(prs):
    slide = _new_slide(prs)
    _title(slide, "C\u1ea5u Tr\u00fac Code \u2014 6 ABAP Include")

    headers = ["Include", "N\u1ed9i dung"]
    rows = [
        [
            "Z_BUG_WS_TOP",
            "Khai b\u00e1o to\u00e0n c\u1ee5c, ki\u1ec3u d\u1eef li\u1ec7u, h\u1eb1ng s\u1ed1 10 tr\u1ea1ng th\u00e1i",
        ],
        ["Z_BUG_WS_F00", "Field catalog cho 5 ALV grid, class LCL_EVENT_HANDLER"],
        [
            "Z_BUG_WS_PBO",
            "Process Before Output \u2014 t\u1ea5t c\u1ea3 9 m\u00e0n h\u00ecnh",
        ],
        ["Z_BUG_WS_PAI", "Process After Input \u2014 to\u00e0n b\u1ed9 fcode handler"],
        [
            "Z_BUG_WS_F01",
            "Business Logic: save, change_status, auto_assign, email, history",
        ],
        ["Z_BUG_WS_F02", "Helper: 10 F4 search help, long text API, SMW0 download"],
    ]
    _table(slide, headers, rows, [0.22, 0.78], BODY_Y, Inches(4.10), fs=Pt(13.5))
    _footer(
        slide,
        "Tri\u1ec3n khai: SE38 \u2192 paste code \u2192 Ctrl+F2 (check) \u2192 Ctrl+F3 (activate)",
        6,
    )


def slide_07_lifecycle(prs):
    slide = _new_slide(prs)
    _title(
        slide, "V\u00f2ng \u0110\u1eddi L\u1ed7i \u2014 10 Tr\u1ea1ng Th\u00e1i (v5.0)"
    )

    headers = ["M\u00e3", "T\u00ean", "Chuy\u1ec3n \u0111\u1ebfn \u0111\u01b0\u1ee3c"]
    rows = [
        ["1", "New", "Assigned (2), Waiting (W)"],
        ["2", "Assigned", "In Progress (3), Rejected (R)"],
        ["3", "In Progress", "Fixed (5), Suspended (4), Rejected (R)"],
        ["4", "Suspended", "Assigned (2)"],
        ["5", "Fixed", "Final Testing (6), Waiting (W)"],
        ["6", "Final Testing", "Resolved (V), In Progress (3)"],
        ["V", "Resolved", "\u2014 k\u1ebft th\u00fac"],
        ["R", "Rejected", "\u2014 k\u1ebft th\u00fac"],
        ["W", "Waiting", "Assigned (2), Final Testing (6)"],
        ["7", "Closed", "\u2014 legacy, t\u01b0\u01a1ng th\u00edch ng\u01b0\u1ee3c"],
    ]
    tbl_h = Inches(4.85)
    _table(slide, headers, rows, [0.06, 0.18, 0.76], BODY_Y, tbl_h, fs=Pt(13))
    _note(
        slide,
        "Quy t\u1eafc: STATUS lu\u00f4n kh\u00f3a tr\u00ean m\u00e0n h\u00ecnh 0300 \u2014 ch\u1ec9 thay \u0111\u1ed5i qua Popup 0370",
        BODY_Y + tbl_h + Inches(0.10),
    )
    _footer(slide, "", 7)


def slide_08_autoassign(prs):
    slide = _new_slide(prs)
    _title(slide, "Engine T\u1ef1 \u0110\u1ed9ng Ph\u00e2n C\u00f4ng")

    _label(
        slide,
        "Giai \u0111o\u1ea1n A \u2014 Khi t\u1ea1o l\u1ed7i (BUG_TYPE = 'C'):",
        BODY_Y,
    )
    code_a_y = BODY_Y + Inches(0.48)
    code_a_h = Inches(1.82)
    _code_block(
        slide,
        [
            "SELECT Dev t\u1eeb ZBUG_USER_PROJEC (role='D', c\u00f9ng project)",
            "JOIN ZBUG_USERS WHERE sap_module = bug.module AND is_active='X'",
            "COUNT active bugs (status IN '2','3','4','6') cho m\u1ed7i dev",
            "\u2192 Ch\u1ecdn dev c\u00f3 workload th\u1ea5p nh\u1ea5t V\u00c0 workload < 5",
            "   T\u00ecm th\u1ea5y  \u2192  STATUS = '2' (Assigned), g\u1eedi email",
            "   Kh\u00f4ng c\u00f3  \u2192  STATUS = 'W' (Waiting), b\u00e1o Manager",
        ],
        code_a_y,
        code_a_h,
    )

    _label(
        slide,
        "Giai \u0111o\u1ea1n B \u2014 Khi chuy\u1ec3n sang Fixed (5):",
        code_a_y + code_a_h + Inches(0.18),
    )
    code_b_y = code_a_y + code_a_h + Inches(0.65)
    code_b_h = Inches(1.60)
    _code_block(
        slide,
        [
            "SELECT Tester t\u1eeb project (role='T', c\u00f9ng SAP module)",
            "COUNT bugs Final Testing \u0111ang active cho m\u1ed7i tester",
            "\u2192 Ch\u1ecdn tester workload th\u1ea5p nh\u1ea5t V\u00c0 workload < 5",
            "   T\u00ecm th\u1ea5y  \u2192  VERIFY_TESTER_ID, STATUS = '6' (Final Testing)",
            "   Kh\u00f4ng c\u00f3  \u2192  STATUS = 'W' (Waiting)",
        ],
        code_b_y,
        code_b_h,
    )

    _footer(slide, "", 8)


def slide_09_access(prs):
    slide = _new_slide(prs)
    _title(slide, "Ki\u1ec3m So\u00e1t Ph\u00e2n Quy\u1ec1n \u2014 Screen Groups")

    headers = [
        "Group",
        "Tr\u01b0\u1eddng b\u1ecb \u1ea3nh h\u01b0\u1edfng",
        "\u0110i\u1ec1u ki\u1ec7n kh\u00f3a",
    ]
    rows = [
        [
            "STS",
            "STATUS",
            "Lu\u00f4n kh\u00f3a \u2014 ch\u1ec9 \u0111\u1ed5i qua Popup 0370",
        ],
        ["BID", "BUG_ID", "Lu\u00f4n kh\u00f3a \u2014 t\u1ef1 t\u1ea1o (BUG0000001)"],
        [
            "PRJ",
            "PROJECT_ID",
            "Kh\u00f3a sau khi \u0111\u1eb7t t\u1eeb context d\u1ef1 \u00e1n",
        ],
        ["FNC", "BUG_TYPE, PRIORITY, SEVERITY", "Kh\u00f3a v\u1edbi role Developer"],
        ["DEV", "DEV_ID, VERIFY_TESTER_ID", "Kh\u00f3a v\u1edbi role Tester"],
        ["TST", "TESTER_ID", "Kh\u00f3a v\u1edbi role Developer"],
        [
            "EDT",
            "T\u1ea5t c\u1ea3 tr\u01b0\u1eddng editable",
            "Kh\u00f3a \u1edf ch\u1ebf \u0111\u1ed9 Display",
        ],
    ]
    _table(slide, headers, rows, [0.08, 0.32, 0.60], BODY_Y, Inches(4.20), fs=Pt(13.5))
    _note(
        slide,
        "C\u01a1 ch\u1ebf: LOOP AT SCREEN ... MODIFY SCREEN trong module PBO",
        Inches(5.48),
    )
    _footer(slide, "", 9)


def slide_10_evidence(prs):
    slide = _new_slide(prs)
    _title(slide, "Qu\u1ea3n L\u00fd B\u1eb1ng Ch\u1ee9ng & Th\u00f4ng B\u00e1o Email")

    _label(
        slide,
        "B\u1eb1ng ch\u1ee9ng \u2014 l\u01b0u trong ZBUG_EVIDENCE (RAWSTRING):",
        BODY_Y,
    )
    ev_h = [
        "Template SMW0",
        "File",
        "Ng\u01b0\u1eddi d\u00f9ng",
        "Th\u1eddi \u0111i\u1ec3m",
    ]
    ev_r = [
        [
            "ZBT_TMPL_01",
            "Bug_report.xlsx",
            "Tester",
            "Khi t\u1ea1o / b\u00e1o c\u00e1o l\u1ed7i",
        ],
        [
            "ZBT_TMPL_02",
            "fix_report.xlsx",
            "Developer",
            "B\u1eaft bu\u1ed9c tr\u01b0\u1edbc khi \u2192 Fixed",
        ],
        [
            "ZBT_TMPL_03",
            "confirm_report.xlsx",
            "Tester",
            "Khi x\u00e1c nh\u1eadn Final Testing",
        ],
    ]
    _table(
        slide,
        ev_h,
        ev_r,
        [0.18, 0.22, 0.16, 0.44],
        BODY_Y + Inches(0.50),
        Inches(1.85),
        fs=Pt(13.5),
    )

    _label(
        slide, "Email t\u1ef1 \u0111\u1ed9ng \u2014 CL_BCS API:", BODY_Y + Inches(2.55)
    )
    em_h = ["S\u1ef1 ki\u1ec7n", "Ng\u01b0\u1eddi nh\u1eadn"]
    em_r = [
        ["CREATE / ASSIGN", "Dev \u0111\u01b0\u1ee3c ph\u00e2n c\u00f4ng + Manager"],
        ["STATUS_CHANGE", "Dev + Tester li\u00ean quan"],
        ["REJECT", "Tester b\u00e1o c\u00e1o l\u1ed7i ban \u0111\u1ea7u"],
    ]
    _table(
        slide,
        em_h,
        em_r,
        [0.30, 0.70],
        BODY_Y + Inches(3.05),
        Inches(1.85),
        fs=Pt(13.5),
    )
    _footer(slide, "", 10)


def slide_11_uat(prs):
    slide = _new_slide(prs)
    _title(slide, "K\u1ebft Qu\u1ea3 UAT V\u00f2ng 1 (11\u201313/04/2026)")

    _label(slide, "T\u1ed5ng quan:", BODY_Y)
    sum_h = [
        "T\u1ed5ng ca",
        "\u0110\u1ea1t",
        "Th\u1ea5t b\u1ea1i",
        "B\u1ecb ch\u1eb7n",
        "T\u1ef7 l\u1ec7 \u0111\u1ea1t",
    ]
    sum_r = [["64", "53", "**11**", "0", "**82.8%**"]]
    _table(
        slide, sum_h, sum_r, [0.2] * 5, BODY_Y + Inches(0.45), Inches(0.92), fs=Pt(14)
    )

    _label(
        slide, "Top 4 l\u1ed7i nghi\u00eam tr\u1ecdng nh\u1ea5t:", BODY_Y + Inches(1.55)
    )
    bug_h = ["ID", "M\u00f4 t\u1ea3", "M\u1ee9c \u0111\u1ed9"]
    bug_r = [
        [
            "UAT-01",
            "Custom Control kh\u00f4ng gi\u1ea3i ph\u00f3ng \u2192 Short dump khi chuy\u1ec3n tab",
            "Critical",
        ],
        [
            "UAT-09",
            "Kh\u00f4ng c\u00f3 ma tr\u1eadn chuy\u1ec3n \u0111\u1ed5i \u2192 cho ph\u00e9p \u0111\u1ea3o ng\u01b0\u1ee3c tr\u1ea1ng th\u00e1i",
            "Critical",
        ],
        [
            "UAT-11",
            "Manager bypass ma tr\u1eadn \u2014 g\u00e1n tr\u1ea1ng th\u00e1i tr\u1ef1c ti\u1ebfp",
            "Critical",
        ],
        ["UAT-08", "Long text bi\u1ebfn m\u1ea5t sau khi save & reopen", "High"],
    ]
    _table(
        slide,
        bug_h,
        bug_r,
        [0.09, 0.76, 0.15],
        BODY_Y + Inches(2.00),
        Inches(2.55),
        fs=Pt(13.5),
    )
    _note(
        slide,
        "T\u1ea5t c\u1ea3 11 l\u1ed7i \u0111\u00e3 ph\u00e2n t\u00edch nguy\u00ean nh\u00e2n v\u00e0 s\u1eeda trong Giai \u0111o\u1ea1n F v5.0 (ho\u00e0n t\u1ea5t 16/04/2026)",
        BODY_Y + Inches(4.72),
    )
    _footer(slide, "", 11)


def slide_12_qc(prs):
    slide = _new_slide(prs)
    _title(slide, "QC Test Plan \u2014 20 Test Suites")

    headers = ["TC", "Suite", "S\u1ed1 ca"]
    rows = [
        [
            "TC-01",
            "Lu\u1ed3ng \u0110i\u1ec1u h\u01b0\u1edbng (t\u1ea5t c\u1ea3 screen transitions)",
            "20",
        ],
        ["TC-08", "Chuy\u1ec3n tr\u1ea1ng th\u00e1i \u2014 10 TT + Popup 0370", "30"],
        [
            "TC-09",
            "Engine T\u1ef1 \u0111\u1ed9ng Ph\u00e2n c\u00f4ng (Phase A + B)",
            "9",
        ],
        ["TC-15", "Ki\u1ec3m so\u00e1t Truy c\u1eadp theo Vai tr\u00f2 (RBAC)", "16"],
        ["TC-11", "Dashboard Metrics (accuracy + realtime refresh)", "12"],
        [
            "TC-10",
            "Bug Search 0210/0220 \u2014 b\u1ed9 l\u1ecdc \u0111a tr\u01b0\u1eddng",
            "15",
        ],
        [
            "TC-19",
            "H\u1ed3i quy \u2014 X\u00e1c minh 11 l\u1ed7i UAT kh\u00f4ng t\u00e1i ph\u00e1t",
            "19",
        ],
        ["TC-20", "Tr\u01b0\u1eddng h\u1ee3p Bi\u00ean & Ranh gi\u1edbi", "20"],
        ["TC-02\u2192TC-18", "C\u00e1c suite c\u00f2n l\u1ea1i", "~69"],
        ["", "**T\u1ed5ng c\u1ed9ng (\u01b0\u1edbc t\u00ednh)**", "**~210**"],
    ]
    _table(slide, headers, rows, [0.14, 0.72, 0.14], BODY_Y, Inches(5.45), fs=Pt(13))
    _footer(
        slide,
        "M\u1ee5c ti\u00eau: \u2265 95% pass  |  0 blocked  |  UAT V\u00f2ng 2 sau deploy",
        12,
    )


def slide_13_timeline(prs):
    slide = _new_slide(prs)
    _title(slide, "Timeline D\u1ef1 \u00c1n \u2014 6 Giai \u0110o\u1ea1n")

    headers = [
        "G\u0110",
        "T\u00ean",
        "K\u1ebft qu\u1ea3 ch\u00ednh",
        "Tr\u1ea1ng th\u00e1i",
    ]
    rows = [
        [
            "A",
            "C\u01a1 s\u1edf CSDL",
            "5 b\u1ea3ng t\u00f9y ch\u1ec9nh, domain, number range",
            "\u2705 Xong",
        ],
        [
            "B",
            "Logic nghi\u1ec7p v\u1ee5",
            "6 Function Module (Create, AutoAssign, Email\u2026)",
            "\u2705 Xong",
        ],
        [
            "C",
            "Giao di\u1ec7n Module Pool",
            "8 m\u00e0n h\u00ecnh, GUI Status, ALV Grid",
            "\u2705 Xong",
        ],
        [
            "D",
            "T\u00ednh n\u0103ng n\u00e2ng cao",
            "SmartForms, Excel upload, F4, long text API",
            "\u2705 Xong",
        ],
        [
            "E",
            "Ki\u1ec3m th\u1eed",
            "QC 140 ca, UAT 64 ca, ph\u00e1t hi\u1ec7n 11 l\u1ed7i",
            "\u2705 Xong",
        ],
        [
            "F",
            "N\u00e2ng cao v5.0",
            "10-state lifecycle, Popup 0370, Dashboard",
            "\u23f3 Deploy ch\u1edd",
        ],
    ]
    _table(
        slide,
        headers,
        rows,
        [0.05, 0.20, 0.59, 0.16],
        BODY_Y,
        Inches(4.00),
        fs=Pt(13.5),
    )
    _footer(
        slide,
        "Code v5.0 ho\u00e0n t\u1ea5t: 16/04/2026  |  Deploy F11\u2013F17: ch\u1edd th\u1ef1c hi\u1ec7n",
        13,
    )


def slide_14_deploy(prs):
    slide = _new_slide(prs)
    _title(slide, "K\u1ebf Ho\u1ea1ch Tri\u1ec3n Khai v5.0 \u2014 7 B\u01b0\u1edbc")

    headers = ["B\u01b0\u1edbc", "C\u00f4ng c\u1ee5", "N\u1ed9i dung"]
    rows = [
        [
            "F11",
            "SE51",
            "T\u1ea1o 4 m\u00e0n h\u00ecnh m\u1edbi: 0410, 0370, 0210, 0220",
        ],
        [
            "F12",
            "SE41",
            "T\u1ea1o 4 GUI Status + Title Bar; c\u1eadp nh\u1eadt STATUS_0200",
        ],
        ["F13", "SE38", "Copy 6 ABAP includes v5.0 \u2192 check \u2192 activate"],
        [
            "F14",
            "SE93",
            "\u0110\u1ed5i m\u00e0n h\u00ecnh ban \u0111\u1ea7u: 0400 \u2192 0410",
        ],
        [
            "F15",
            "SE11",
            "T\u1ea1o b\u1ea3ng ZBUG_EVIDENCE (RAWSTRING, 11 tr\u01b0\u1eddng)",
        ],
        ["F16", "SE38", "Migration: status='6' \u2192 status='V' + COMMIT WORK"],
        ["F17", "SMW0", "Upload 3 file m\u1eabu: ZBT_TMPL_01 / 02 / 03"],
    ]
    _table(slide, headers, rows, [0.08, 0.09, 0.83], BODY_Y, Inches(4.25), fs=Pt(13.5))
    _footer(
        slide,
        "X\u00e1c minh: /nZBUG_WS \u2192 M\u00e0n h\u00ecnh 0410 xu\u1ea5t hi\u1ec7n \u0111\u1ea7u ti\u00ean",
        14,
    )


def slide_15_summary(prs):
    slide = _new_slide(prs)
    _title(slide, "T\u1ed5ng K\u1ebft")

    tf = _txb(slide, LEFT, BODY_Y, CW, BODY_H)

    p0 = tf.paragraphs[0]
    _run(
        p0,
        "\u0110\u00e3 \u0111\u1ea1t \u0111\u01b0\u1ee3c:",
        size=Pt(16.5),
        color=CYAN,
        bold=True,
    )

    done = [
        "H\u1ec7 th\u1ed1ng Bug Tracking ch\u1ea1y native tr\u00ean SAP \u2014 kh\u00f4ng c\u1ea7n license b\u1ed5 sung",
        "V\u00f2ng \u0111\u1eddi 10 tr\u1ea1ng th\u00e1i \u0111\u01b0\u1ee3c th\u1ef1c thi theo vai tr\u00f2 qua Popup 0370",
        "Auto-assign 2 giai \u0111o\u1ea1n d\u1ef1a tr\u00ean SAP module + workload",
        "Audit log b\u1ea5t bi\u1ebfn \u2014 m\u1ecdi thay \u0111\u1ed5i ghi v\u00e0o ZBUG_HISTORY",
        "Dashboard realtime tr\u00ean m\u00e0n h\u00ecnh 0200",
        "11 l\u1ed7i UAT \u2192 ph\u00e2n t\u00edch v\u00e0 s\u1eeda ho\u00e0n to\u00e0n trong v5.0",
    ]
    for item in done:
        p = tf.add_paragraph()
        _run(p, f"  \u2022  {item}", size=Pt(15), color=LIGHT)

    p_sp = tf.add_paragraph()
    _run(p_sp, "", size=Pt(5))

    p_lbl = tf.add_paragraph()
    _run(p_lbl, "C\u00f2n l\u1ea1i:", size=Pt(16.5), color=PINK, bold=True)

    todo = [
        "Tri\u1ec3n khai v5.0 l\u00ean SAP (F11\u2013F17)",
        "Ch\u1ea1y QC Test \u0111\u1ea7y \u0111\u1ee7 20 suites (~210 ca)",
        "UAT V\u00f2ng 2 \u2192 3 th\u00e0nh vi\u00ean k\u00fd x\u00e1c nh\u1eadn ch\u1ea5p thu\u1eadn",
    ]
    for item in todo:
        p = tf.add_paragraph()
        _run(p, f"  \u2022  {item}", size=Pt(15), color=LIGHT)

    _footer(slide, "", 15)


def slide_16_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_DARK

    _rect(slide, 0, 0, Inches(0.12), int(H * 0.55), CYAN)
    _rect(slide, 0, int(H * 0.55), Inches(0.12), int(H * 0.45), PINK)
    _rect(slide, 0, H - Emu(30000), W, Emu(30000), CYAN)

    tf1 = _txb(slide, Inches(0.90), Inches(2.20), Inches(11.5), Inches(1.10))
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    _run(
        p1,
        "C\u1ea3m \u01a0n & H\u1ecfi \u0110\u00e1p",
        size=Pt(44),
        color=CYAN,
        bold=True,
    )

    _rect(slide, Inches(3.0), Inches(3.45), Inches(7.5), Emu(38000), PINK)

    tf2 = _txb(slide, Inches(0.90), Inches(3.65), Inches(11.5), Inches(0.78))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    _run(
        p2,
        "ZBUG_WS v5.0 \u2014 SAP Bug Tracking Management System",
        size=Pt(18),
        color=LIGHT,
    )

    tf3 = _txb(slide, Inches(0.90), Inches(4.58), Inches(11.5), Inches(0.55))
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    _run(
        p3,
        "Nh\u00f3m ZBUG  \u2502  FPT University  \u2502  Th\u00e1ng 4 n\u0103m 2026",
        size=Pt(15),
        color=MUTED,
    )

    tf4 = _txb(slide, W - Inches(1.25), FOOT_Y, Inches(1.15), FOOT_H)
    p4 = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.RIGHT
    _run(p4, f"{TOTAL} / {TOTAL}", size=Pt(11), color=MUTED)


# ── Main ───────────────────────────────────────────────────────────────────────


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_01_title(prs)
    slide_02_team(prs)
    slide_03_context(prs)
    slide_04_comparison(prs)
    slide_05_screens(prs)
    slide_06_code(prs)
    slide_07_lifecycle(prs)
    slide_08_autoassign(prs)
    slide_09_access(prs)
    slide_10_evidence(prs)
    slide_11_uat(prs)
    slide_12_qc(prs)
    slide_13_timeline(prs)
    slide_14_deploy(prs)
    slide_15_summary(prs)
    slide_16_thanks(prs)

    out = os.path.join(
        os.path.dirname(os.path.abspath(__file__)), "output", "presentation_py.pptx"
    )
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print(f"Saved  : {out}")
    print(f"Slides : {len(prs.slides)}")


if __name__ == "__main__":
    main()
